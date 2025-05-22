import streamlit as st
from pptx import Presentation
from sentence_transformers import SentenceTransformer, util
import gdown
import re
import os

st.title("📊 AI Q&A from PowerPoint (Google Drive)")

model = SentenceTransformer("all-mpnet-base-v2")

def download_file_from_drive(gdrive_url, save_path):
    match = re.search(r'/d/([a-zA-Z0-9_-]+)', gdrive_url)
    if not match:
        raise ValueError("Invalid Google Drive URL")
    file_id = match.group(1)
    gdown.download(f"https://drive.google.com/uc?id={file_id}", save_path, quiet=False)

def extract_slide_knowledge(pptx_path):
    prs = Presentation(pptx_path)
    slide_knowledge = []
    for slide in prs.slides:
        title = slide.shapes.title.text if slide.shapes.title else ""
        content = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text != title:
                content += shape.text.strip() + "\n"
        if title or content:
            slide_knowledge.append(f"{title.strip()}\n{content.strip()}")
    return slide_knowledge

# Streamlit input UI
gdrive_url = st.text_input("🔗 Paste your Google Drive PPTX Link here")
question = st.text_input("❓ Enter your question")

if st.button("Get Answer"):
    if not gdrive_url or not question:
        st.warning("Please provide both a link and a question.")
    else:
        try:
            save_path = "presentation.pptx"
            download_file_from_drive(gdrive_url, save_path)
            slides = extract_slide_knowledge(save_path)
            if not slides:
                st.error("No content found in slides.")
            else:
                slide_embeddings = model.encode(slides, convert_to_tensor=True)
                question_embedding = model.encode(question, convert_to_tensor=True)
                scores = util.cos_sim(question_embedding, slide_embeddings)
                best_idx = scores.argmax().item()
                confidence = scores[0][best_idx].item()
                best_slide = slides[best_idx]

                st.success("✅ Best matching answer found:")
                st.markdown(f"""**Slide Content:**  
{best_slide}""")
                st.markdown(f"**Confidence Score:** {confidence:.2f}")

        except Exception as e:
            st.error(f"Error: {str(e)}")
